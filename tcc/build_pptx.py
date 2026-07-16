"""Gera apresentação PPTX da defesa do TCC com base no template oficial.

Usa o slide 0 do template como capa (preserva logo) e cria slides adicionais
com o layout "Título e Conteúdo".
"""
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


TEMPLATE = Path("/home/akirem/Downloads/Template Apresentação.pptx")
OUT = Path("/home/akirem/Downloads/TCC Apresentação.pptx")

# Conteúdo do TCC
TITULO = "Geração automática de texto com LLMs para atendimento ao cliente ou resumo de documentos"
ALUNO = "João Pedro Azevedo Veras dos Reis"
ORIENTADORA = "Prof.ª Gabrielle Lombardi"
INSTITUICAO = "MBA USP/Esalq — Data Science and Analytics"
ANO = "2025"

GRAFICO_PATH = Path("/home/akirem/Documentos/LLM-USP/outputs/grafico_loss_ptbr.png")


def set_title(slide, text: str, size: int = 32):
    """Define o título preservando posição do placeholder."""
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == 0 or "Título" in (sh.name or ""):
            tf = sh.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = True
            return
    # fallback: primeiro shape com text_frame
    for sh in slide.shapes:
        if sh.has_text_frame and "Título" in (sh.name or ""):
            sh.text_frame.text = text
            for r in sh.text_frame.paragraphs[0].runs:
                r.font.size = Pt(size); r.font.bold = True
            return


def set_body_bullets(slide, lines: list[str], size: int = 18):
    """Preenche o placeholder de conteúdo com bullets (uma linha por bullet)."""
    body_ph = None
    for sh in slide.placeholders:
        if sh.placeholder_format.idx != 0:
            body_ph = sh
            break
    if body_ph is None:
        # Cria caixa de texto
        from pptx.util import Inches
        body_ph = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = body_ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = 0
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)


def add_image_slide(prs, layout, title: str, image_path: Path, caption: str = "", img_top_in: float = 1.3, img_height_in: float = 5.5):
    """Slide com título + imagem centralizada + legenda opcional."""
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    # Remove placeholder de body se houver (vamos colocar imagem)
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0:
            sp = ph._element
            sp.getparent().remove(sp)
    # Adiciona imagem
    if image_path.exists():
        slide_w_in = Emu(prs.slide_width).inches
        img_w = 9.0   # largura fixa em inches
        left = Inches((slide_w_in - img_w) / 2)
        slide.shapes.add_picture(str(image_path), left, Inches(img_top_in),
                                  width=Inches(img_w), height=Inches(img_height_in))
    if caption:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(img_top_in + img_height_in + 0.1),
                                       Inches(12), Inches(0.6))
        tf = tb.text_frame
        tf.text = caption
        for r in tf.paragraphs[0].runs:
            r.font.size = Pt(14); r.font.italic = True
    return slide


def add_table_slide(prs, layout, title: str, rows: list[list[str]], col_widths_in: list[float]):
    """Slide com título + tabela."""
    from pptx.util import Inches
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0:
            sp = ph._element
            sp.getparent().remove(sp)
    n_rows = len(rows)
    n_cols = len(rows[0])
    total_w = sum(col_widths_in)
    slide_w_in = Emu(prs.slide_width).inches
    left = Inches((slide_w_in - total_w) / 2)
    top = Inches(1.5)
    height = Inches(0.5 * n_rows + 0.5)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top,
                                   Inches(total_w), height).table
    for ci, w in enumerate(col_widths_in):
        tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14 if ri > 0 else 15)
                    run.font.bold = (ri == 0)
                if ri == 0:
                    para.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, layout, title: str, bullets: list[str], size: int = 18):
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    set_body_bullets(slide, bullets, size=size)
    return slide


def main():
    prs = Presentation(str(TEMPLATE))
    title_layout = prs.slide_masters[0].slide_layouts[0]   # Slide de Título
    content_layout = prs.slide_masters[0].slide_layouts[1] # Título e Conteúdo
    section_layout = prs.slide_masters[0].slide_layouts[2] # Cabeçalho da Seção

    # ============================================================
    # SLIDE 1 — CAPA (edita o slide 0 do template)
    # ============================================================
    capa = prs.slides[0]
    for ph in capa.placeholders:
        if ph.placeholder_format.idx == 0:  # Título
            ph.text_frame.text = TITULO
            for r in ph.text_frame.paragraphs[0].runs:
                r.font.size = Pt(28); r.font.bold = True
        else:  # Subtítulo
            ph.text_frame.text = ""
            paras = [
                ALUNO,
                ORIENTADORA,
                f"{INSTITUICAO} · {ANO}",
            ]
            for i, line in enumerate(paras):
                p = ph.text_frame.paragraphs[0] if i == 0 else ph.text_frame.add_paragraph()
                r = p.add_run()
                r.text = line
                r.font.size = Pt(18)
                r.font.bold = (i == 0)

    # SLIDE 2 — Agenda (reutiliza o slide 1 do template em vez de deletar)
    # O slide 1 do template não tem placeholder de título → adiciona um TextBox.
    agenda = prs.slides[1]
    title_box = agenda.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Roteiro da Apresentação"
    for r in title_tf.paragraphs[0].runs:
        r.font.size = Pt(32); r.font.bold = True
    set_body_bullets(agenda, [
        "1. Contexto e motivação",
        "2. Objetivos do trabalho",
        "3. Materiais e métodos",
        "4. Arquitetura e pipeline de treinamento",
        "5. Resultados quantitativos e qualitativos",
        "6. Discussão à luz das leis de escala",
        "7. Limitações e trabalhos futuros",
        "8. Conclusões",
    ])

    # ============================================================
    # SLIDE 3 — Considerações Iniciais
    # ============================================================
    add_content_slide(prs, content_layout, "Considerações Iniciais", [
        "LLMs tornaram-se centrais no PLN após a arquitetura Transformer (Vaswani et al., 2017).",
        "Leis de escala (Kaplan et al., 2020; Hoffmann et al., 2022) relacionam parâmetros, dados e desempenho.",
        "Alinhamento via SFT e RLHF é necessário para uso seguro (Ouyang et al., 2022).",
        "No Brasil, há carência de modelos nativos em PT-BR e dependência de provedores internacionais.",
    ])

    # ============================================================
    # SLIDE 4 — Problema e motivação
    # ============================================================
    add_content_slide(prs, content_layout, "Problema e Motivação", [
        "Modelos proprietários (ex.: GPT-4) impõem barreiras à adoção corporativa segura no Brasil.",
        "Riscos: conformidade com a LGPD, viés cultural, latência e custo indexado ao dólar.",
        "Hipótese: é viável desenvolver um SLM nativo em PT-BR em hardware acadêmico?",
        "Foco em aplicações futuras de atendimento ao cliente e resumo de documentos.",
    ])

    # ============================================================
    # SLIDE 5 — Objetivos
    # ============================================================
    add_content_slide(prs, content_layout, "Objetivos", [
        "Objetivo geral: desenvolver, treinar e avaliar um SLM decoder-only em PT-BR em 6 GB de VRAM.",
        "(i) Compilar e normalizar corpus multi-domínio em português.",
        "(ii) Definir arquitetura Transformer decoder-only sob restrição de hardware.",
        "(iii) Executar pré-treino causal com monitoramento contínuo de loss e acurácia.",
        "(iv) Avaliar qualitativa e quantitativamente o modelo resultante.",
        "SFT e alinhamento por preferências permanecem como trabalhos futuros.",
    ])

    # ============================================================
    # SLIDE 6 — Materiais e Hardware
    # ============================================================
    add_content_slide(prs, content_layout, "Materiais e Hardware", [
        "GPU NVIDIA GeForce RTX 3050 Laptop (6 GB de VRAM, CUDA).",
        "CPU Intel Core i5-13420H · 7 GB de memória RAM.",
        "Python 3.12.3 · PyTorch 2.11 · Hugging Face datasets.",
        "Tokenização: BPE do GPT-2 (50.257 tokens).",
        "Pipeline 100% local, sem dependência de provedores externos.",
    ])

    # ============================================================
    # SLIDE 7 — Corpus Multi-Domínio
    # ============================================================
    add_content_slide(prs, content_layout, "Corpus Multi-Domínio", [
        "Jurídico: Constituição Federal (3.575 linhas).",
        "Literário clássico: obras de domínio público (Machado de Assis).",
        "Ficção contemporânea: amostras em uso acadêmico não-comercial.",
        "V7: incorporação de pares pergunta/resposta traduzidos (Dolly 15k via LibreTranslate + Alpaca via OPUS-MT).",
        "Volume final V7: ~37 MB, ~2,5 milhões de tokens.",
        "Split treino/validação: corte sequencial 90 % / 10 %.",
    ])

    # ============================================================
    # SLIDE 8 — Arquitetura (Tabela)
    # ============================================================
    add_table_slide(prs, content_layout, "Arquitetura do Modelo (SLM)", [
        ["Hiperparâmetro", "Valor"],
        ["Arquitetura", "Transformer Decoder-Only (GPT-2)"],
        ["Parâmetros totais", "52,93 milhões"],
        ["Camadas (n_layer)", "8"],
        ["Cabeças de atenção (n_head)", "6"],
        ["Dimensão do embedding (n_embd)", "384"],
        ["Janela de contexto (block_size)", "256 tokens"],
        ["Batch size efetivo", "16 (1 × grad_accum 16)"],
        ["Learning rate inicial", "3e-4 (CosineAnnealing → 1e-5)"],
        ["Dropout", "0,1"],
    ], col_widths_in=[5.0, 5.5])

    # ============================================================
    # SLIDE 9 — Pipeline de Treinamento
    # ============================================================
    add_content_slide(prs, content_layout, "Pipeline de Treinamento", [
        "Função de perda: entropia cruzada (cross-entropy) por token.",
        "Otimização: AdamW (β₁=0,9; β₂=0,999; weight decay 1e-2).",
        "Scheduler: CosineAnnealingLR (3e-4 → 1e-5).",
        "Gradient clipping com norma máxima 1,0.",
        "Avaliação intermediária a cada 200 iterações.",
        "Critério de parada: estabilização empírica da perda de validação.",
    ])

    # ============================================================
    # SLIDE 10 — Curva de Loss (gráfico)
    # ============================================================
    add_image_slide(prs, content_layout, "Resultados: Evolução da Loss",
                     GRAFICO_PATH,
                     caption="Figura 1. Loss de treino e validação (Cross-Entropy) ao longo das iterações.",
                     img_top_in=1.3, img_height_in=5.2)

    # ============================================================
    # SLIDE 11 — Métricas (Tabela)
    # ============================================================
    add_table_slide(prs, content_layout, "Resultados: Métricas por Iteração", [
        ["Iteração", "Train Loss", "Val Loss", "Acurácia (Val)"],
        ["200 (V6)", "4,144", "4,286", "26,2 %"],
        ["800 (V6)", "3,363", "3,008", "35,6 %"],
        ["1.400 (V6)", "1,044", "2,145", "63,5 %"],
        ["3.000 (V7)", "1,835", "2,029", "58,8 %"],
        ["5.000 (V7)", "1,775", "2,003", "59,4 %"],
    ], col_widths_in=[3.0, 2.5, 2.5, 2.5])

    # ============================================================
    # SLIDE 12 — Amostras Qualitativas
    # ============================================================
    add_content_slide(prs, content_layout, "Resultados: Amostras Qualitativas (V7)", [
        "Prompt: \"Era uma vez um menino que\"",
        "→ \"...me perguntava ao dia. ia ajudá-la a irmã em si mesma. Como ia lembrar-se...\"",
        " ",
        "Prompt: \"A velha senhora olhou pela janela e\"",
        "→ \"...sendo ingerindo seu cabelo, meu corpo está encarando o botão...\"",
        " ",
        "Observações: respeita ortografia/acentuação, reproduz estruturas de diálogo,",
        "mas apresenta repetição de muletas e perda de coerência semântica após ~15 tokens.",
    ], size=15)

    # ============================================================
    # SLIDE 13 — Discussão (Leis de Escala)
    # ============================================================
    add_content_slide(prs, content_layout, "Discussão: Leis de Escala", [
        "V6 vs V7 — aparente regressão da Train Loss não indica piora do aprendizado.",
        "V6 (corpus pequeno, ~6 MB): memorização → gap train/val ≈ 1,10 (overfitting).",
        "V7 (corpus ampliado, ~37 MB): generalização → gap train/val ≈ 0,22.",
        "Modelo aprende sintaxe, conjugação e diálogo, mas não internaliza fatos específicos.",
        "Tentativa de SFT direcionado com 162 pares Q&A curados em PT-BR: não houve incorporação factual.",
        "Achado consistente com Kaplan et al. (2020): escala paramétrica é o gargalo, não dados.",
    ], size=16)

    # ============================================================
    # SLIDE 14 — Limitações e Trabalhos Futuros
    # ============================================================
    add_content_slide(prs, content_layout, "Limitações e Trabalhos Futuros", [
        "Limitações reconhecidas:",
        "  • Corpus modesto (2,5 milhões de tokens) e janela de contexto reduzida (256).",
        "  • Tokenizador GPT-2 herdado, sub-ótimo para PT-BR.",
        "  • Avaliação intrínseca (loss/acurácia); sem benchmarks downstream.",
        " ",
        "Trabalhos futuros:",
        "  • Fine-Tuning Supervisionado (SFT) e alinhamento por preferências (RLHF).",
        "  • Retrieval-Augmented Generation (RAG) com Wikipedia-PT.",
        "  • Modelos maiores (Tucano-630m/1b1) com LoRA e quantização 4-bit.",
        "  • Avaliação via ROUGE-L, BERTScore e julgamento humano.",
    ], size=15)

    # ============================================================
    # SLIDE 15 — Conclusões
    # ============================================================
    add_content_slide(prs, content_layout, "Conclusões", [
        "Demonstrou-se a viabilidade técnica de treinar um SLM em PT-BR em hardware acadêmico (6 GB VRAM).",
        "Pipeline reproduzível: curadoria → tokenização → arquitetura decoder-only → pré-treino causal.",
        "Loss de validação reduzida de 4,29 → 2,00; acurácia de 26,2 % → 59,4 %.",
        "Validação experimental das leis de escala: escala paramétrica é gargalo factual.",
        "Próximos passos: SFT, RAG e expansão de modelo são caminhos claros para evolução.",
    ])

    # ============================================================
    # SLIDE 16 — Obrigado
    # ============================================================
    slide_fim = prs.slides.add_slide(section_layout)
    set_title(slide_fim, "Obrigado!", size=54)
    # Subtítulo
    body_ph = None
    for ph in slide_fim.placeholders:
        if ph.placeholder_format.idx != 0:
            body_ph = ph; break
    if body_ph is not None:
        body_ph.text_frame.text = (
            f"{ALUNO}\n"
            f"{ORIENTADORA}\n"
            f"{INSTITUICAO} · {ANO}"
        )
        for p in body_ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(20)

    prs.save(str(OUT))
    print(f"[ok] apresentação salva: {OUT}")
    print(f"[ok] tamanho: {OUT.stat().st_size/1e6:.2f} MB")
    print(f"[ok] {len(prs.slides)} slides gerados")


if __name__ == "__main__":
    main()
